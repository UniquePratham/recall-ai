from PyPDF2 import PdfReader
import docx
import requests
from bs4 import BeautifulSoup
import speech_recognition as sr
from pydub import AudioSegment
from io import BytesIO
# python-pptx imported conditionally inside process_powerpoint function
from utils import extract_text_from_image, summarize_text, generate_embedding, store_in_vector_db, query_knowledge_base


async def process_document(doc, username):
    file = await doc.get_file()
    file_content = await file.download_as_bytearray()

    if doc.file_name.endswith('.pdf'):
        return await process_pdf(BytesIO(file_content), username)
    elif doc.file_name.endswith('.docx'):
        return await process_docx(BytesIO(file_content), username)
    elif doc.file_name.endswith(('.html', '.htm')):
        return await process_html(file_content, username, doc.file_name)
    elif doc.file_name.endswith('.txt'):
        return await process_txt(file_content, username)
    elif doc.file_name.endswith('.md'):
        return await process_markdown(file_content, username)
    elif doc.file_name.endswith(('.ppt', '.pptx')):
        return await process_powerpoint(BytesIO(file_content), username, doc.file_name)
    else:
        return "Unsupported document type"


async def process_pdf(file_obj, username):
    reader = PdfReader(file_obj)
    if len(reader.pages) > 30:
        return "PDF is too long (over 30 pages)"

    text = ""
    for page in reader.pages:
        text += page.extract_text()

    if not text:
        # If no text extracted, it might be a scanned PDF
        return await process_scanned_pdf(file_obj, username)

    summary = await summarize_text(text)
    embedding = await generate_embedding(summary)
    await store_in_vector_db(embedding, summary, username)
    return "PDF processed and stored"


async def process_scanned_pdf(file_obj, username):
    # Use OCR or GPT-4 Vision here
    text = await extract_text_from_image(file_obj)
    summary = await summarize_text(text)
    embedding = await generate_embedding(summary)
    await store_in_vector_db(embedding, summary, username)
    return "Scanned PDF processed and stored"


async def process_docx(file_obj, username):
    doc = docx.Document(file_obj)

    # Count total words
    total_words = sum(len(paragraph.text.split())
                      for paragraph in doc.paragraphs)

    # Estimate pages (assuming average 500 words per page)
    estimated_pages = total_words // 400

    if estimated_pages > 30:
        return f"DOCX is too long (estimated {estimated_pages} pages, over 30 page limit)"

    text = "\n".join([para.text for para in doc.paragraphs])
    summary = await summarize_text(text)
    embedding = await generate_embedding(summary)
    await store_in_vector_db(embedding, summary, username, metadata={"content_type": "docx", "pages": estimated_pages})
    return f"DOCX processed and stored (estimated {estimated_pages} pages)"


async def process_html(file_content, username, filename):
    """Process HTML files and extract text content"""
    try:
        # Decode the file content
        html_content = file_content.decode('utf-8', errors='ignore')

        # Parse with BeautifulSoup
        soup = BeautifulSoup(html_content, 'html.parser')

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Extract title
        title = soup.title.string if soup.title else filename

        # Get text content
        text = soup.get_text()

        # Clean up whitespace
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip()
                  for line in lines for phrase in line.split("  "))
        clean_text = ' '.join(chunk for chunk in chunks if chunk)

        if len(clean_text) > 10000:
            clean_text = clean_text[:10000] + "..."

        # Create content with title for better context
        content = f"HTML Document: {title}\n\nContent: {clean_text}"

        if len(content.strip()) < 50:
            return "HTML file appears to be empty or contains no readable text"

        summary = await summarize_text(content)
        embedding = await generate_embedding(summary)
        await store_in_vector_db(embedding, summary, username, metadata={"content_type": "html", "title": title})

        return f"HTML file '{filename}' processed and stored"

    except Exception as e:
        return f"Error processing HTML file: {str(e)}"


async def process_txt(file_content, username):
    """Process plain text files"""
    try:
        # Decode the file content
        text = file_content.decode('utf-8', errors='ignore').strip()

        if not text:
            return "Text file is empty"

        # Check file size (approximate)
        word_count = len(text.split())
        if word_count > 15000:  # Roughly 30 pages
            return f"Text file is too long ({word_count} words, over 15000 word limit)"

        # For large texts, summarize. For smaller ones, store directly
        if len(text) > 2000:
            summary = await summarize_text(text)
            content_to_store = summary
        else:
            content_to_store = text

        embedding = await generate_embedding(content_to_store)
        await store_in_vector_db(embedding, content_to_store, username, metadata={"content_type": "txt", "word_count": word_count})

        return f"Text file processed and stored ({word_count} words)"

    except Exception as e:
        return f"Error processing text file: {str(e)}"


async def process_markdown(file_content, username):
    """Process Markdown files"""
    try:
        # Decode the file content
        markdown_text = file_content.decode('utf-8', errors='ignore').strip()

        if not markdown_text:
            return "Markdown file is empty"

        # Check file size
        word_count = len(markdown_text.split())
        if word_count > 15000:
            return f"Markdown file is too long ({word_count} words, over 15000 word limit)"

        # For large markdown files, summarize. For smaller ones, store directly
        if len(markdown_text) > 2000:
            summary = await summarize_text(markdown_text)
            content_to_store = f"Markdown Document Summary: {summary}"
        else:
            content_to_store = f"Markdown Document: {markdown_text}"

        embedding = await generate_embedding(content_to_store)
        await store_in_vector_db(embedding, content_to_store, username, metadata={"content_type": "markdown", "word_count": word_count})

        return f"Markdown file processed and stored ({word_count} words)"

    except Exception as e:
        return f"Error processing Markdown file: {str(e)}"


async def process_powerpoint(file_obj, username, filename):
    """Process PowerPoint files (.ppt, .pptx)"""
    try:
        # Import python-pptx here to avoid import errors if not installed
        try:
            from pptx import Presentation
        except ImportError:
            return "PowerPoint processing requires python-pptx library. Please install it with: pip install python-pptx"

        # Load the presentation
        prs = Presentation(file_obj)
        slide_count = len(prs.slides)

        if slide_count > 50:
            return f"PowerPoint is too long ({slide_count} slides, over 50 slide limit)"

        # Extract text from all slides
        all_text = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text.strip())

            if slide_text:
                all_text.append(f"Slide {i}: {' | '.join(slide_text)}")

        if not all_text:
            return "No text content found in PowerPoint file"

        text = "\n\n".join(all_text)
        word_count = len(text.split())

        # Create comprehensive content
        content = f"PowerPoint Presentation: {filename}\nSlides: {slide_count}\nContent:\n{text}"

        # Summarize if too long
        if len(content) > 3000:
            summary = await summarize_text(content)
            content_to_store = f"PowerPoint: {filename} ({slide_count} slides)\nSummary: {summary}"
        else:
            content_to_store = content

        embedding = await generate_embedding(content_to_store)
        await store_in_vector_db(embedding, content_to_store, username,
                                 metadata={"content_type": "powerpoint", "slides": slide_count, "filename": filename})

        return f"PowerPoint '{filename}' processed and stored ({slide_count} slides, {word_count} words)"

    except Exception as e:
        return f"Error processing PowerPoint file: {str(e)}"


async def process_photo(photo, username):
    file = await photo.get_file()
    summary = await extract_text_from_image(file.file_path)
    embedding = await generate_embedding(summary)
    await store_in_vector_db(embedding, summary, username)
    return f"Photo analyzed and stored: {summary}"


async def process_audio(audio, username):
    try:
        file = await audio.get_file()
        file_content = await file.download_as_bytearray()

        # Convert MP3 to WAV
        audio_segment = AudioSegment.from_mp3(BytesIO(file_content))
        wav_io = BytesIO()
        audio_segment.export(wav_io, format="wav")
        wav_io.seek(0)

        recognizer = sr.Recognizer()
        with sr.AudioFile(wav_io) as source:
            audio_data = recognizer.record(source)
            text = recognizer.recognize_google(audio_data)

        summary = await summarize_text(text)
        embedding = await generate_embedding(summary)
        await store_in_vector_db(embedding, summary, username)
        return f"Audio transcribed and stored. Transcription: {text}"
    except sr.UnknownValueError:
        return "Speech Recognition could not understand the audio"
    except sr.RequestError as e:
        return f"Could not request results from Speech Recognition service; {e}"
    except Exception as e:
        return f"Error processing audio: {str(e)}"


async def process_text(text, username):
    if text.lower().startswith(("what", "how", "why", "when", "where", "who")):
        return await query_knowledge_base(text, username)
    embedding = await generate_embedding(text)
    await store_in_vector_db(embedding, text, username)
    return "Text processed and stored"


async def process_url(url, username):
    try:
        # Add headers to appear more like a regular browser
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36',
            'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,image/webp,*/*;q=0.8',
            'Accept-Language': 'en-US,en;q=0.5',
            'Accept-Encoding': 'gzip, deflate',
            'Connection': 'keep-alive',
            'Upgrade-Insecure-Requests': '1',
        }

        response = requests.get(url, headers=headers, timeout=10)
        response.raise_for_status()

        soup = BeautifulSoup(response.text, 'html.parser')

        # Extract title
        title = soup.title.string.strip() if soup.title and soup.title.string else "No title"

        # Remove script and style elements
        for script in soup(["script", "style"]):
            script.decompose()

        # Get text content
        text = soup.get_text()

        # Clean up the text
        lines = (line.strip() for line in text.splitlines())
        chunks = (phrase.strip()
                  for line in lines for phrase in line.split("  "))
        text = ' '.join(chunk for chunk in chunks if chunk)

        # Limit text length to avoid overwhelming the system
        if len(text) > 5000:
            text = text[:5000] + "..."

        # Create comprehensive content for storage - INCLUDE URL AND TITLE PROMINENTLY
        if text.strip():
            # Include URL and description in the content for better searchability
            url_description = ""
            if any(keyword in url.lower() for keyword in ["ai", "coding", "builder", "dev", "emergent", "v0", "mgx", "gamma"]):
                url_description = " - AI Coding Website and Project Builder"
            elif "cricket" in url.lower() or "sports" in url.lower():
                url_description = " - Sports Streaming Website"

            content = f"Website: {url}{url_description}\nTitle: {title}\nContent: {text}"
            summary = await summarize_text(content, max_length=800)

            # Create storage text that prominently features the URL and description
            storage_text = f"URL: {url}{url_description}\nTitle: {title}\nSummary: {summary}"
        else:
            # Handle cases where content extraction fails (404, bot protection, etc.)
            url_description = " - AI Coding Website and Project Builder" if any(keyword in url.lower(
            ) for keyword in ["emergent", "v0", "mgx", "gamma", "ai", "coding", "builder", "dev"]) else ""
            storage_text = f"URL: {url}{url_description}\nTitle: {title}\nNote: Content could not be extracted (possibly due to 404 error or bot protection)"
            summary = storage_text

        # Store with URL prominently in the embedding text for better search
        embedding_text = f"Website: {url}{url_description} - {title}. {summary}"
        embedding = await generate_embedding(embedding_text)

        await store_in_vector_db(
            embedding,
            storage_text,
            username,
            metadata={"url": url, "title": title, "content_type": "url"}
        )

        return "URL content extracted and stored"

    except requests.exceptions.RequestException as e:
        # Handle network errors, 404s, timeouts, etc.
        url_description = " - AI Coding Website and Project Builder" if any(keyword in url.lower(
        ) for keyword in ["emergent", "v0", "mgx", "gamma", "ai", "coding", "builder", "dev"]) else ""
        error_info = f"URL: {url}{url_description}\nNote: Could not access website ({str(e)}). This might be due to bot protection, 404 error, or network issues."

        # Still store the URL information even if we can't access it
        embedding = await generate_embedding(f"Website: {url}{url_description} - Access error: {str(e)}")
        await store_in_vector_db(
            embedding,
            error_info,
            username,
            metadata={"url": url, "content_type": "url", "error": str(e)}
        )

        return f"URL noted but content inaccessible: {str(e)}"

    except Exception as e:
        # Handle other errors
        return f"Error processing URL: {str(e)}"
